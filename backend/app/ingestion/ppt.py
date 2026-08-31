"""
PowerPoint ingestion module for VisionRAG-X.

Extracts text, images and speaker notes from PPTX files using python-pptx.
"""

import asyncio
import io
import logging
from pathlib import Path
from typing import List

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

logger = logging.getLogger(__name__)


class PPTIngestionError(Exception):
    """Raised when PPTX ingestion fails."""


class PPTIngester:
    """
    Extracts structured content from PowerPoint (.pptx) files.

    Iterates over slides, collecting text from all text-frame shapes,
    exporting embedded images, and capturing speaker notes.
    """

    # ------------------------------------------------------------------ #
    # Extraction
    # ------------------------------------------------------------------ #

    async def extract(self, ppt_path: Path, output_dir: Path) -> List[dict]:
        """
        Extract text, images, and notes from every slide in a PPTX.

        Parameters
        ----------
        ppt_path:
            Path to the .pptx file.
        output_dir:
            Root directory for extracted assets. Images are saved under
            ``output_dir/images/``.

        Returns
        -------
        List of slide dicts::

            [
                {
                    "slide_num": int,          # 1-indexed
                    "title": str,
                    "text": str,
                    "images": [{"image_path": str}, ...],
                    "notes": str,
                },
                ...
            ]
        """
        ppt_path = Path(ppt_path)
        output_dir = Path(output_dir)

        if not ppt_path.exists():
            raise PPTIngestionError(f"PPTX file not found: {ppt_path}")

        images_dir = output_dir / "images"
        images_dir.mkdir(parents=True, exist_ok=True)

        def _extract():
            try:
                prs = Presentation(str(ppt_path))
            except Exception as exc:
                raise PPTIngestionError(
                    f"python-pptx could not open {ppt_path}: {exc}"
                ) from exc

            slides_data = []
            for slide_idx, slide in enumerate(prs.slides):
                slide_num = slide_idx + 1

                title_text = ""
                text_parts: List[str] = []
                image_records: List[dict] = []
                img_counter = 0

                for shape in slide.shapes:
                    # ---- Text frames ----
                    if shape.has_text_frame:
                        shape_text = "\n".join(
                            para.text
                            for para in shape.text_frame.paragraphs
                            if para.text.strip()
                        )
                        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            # Try to identify the title placeholder
                            try:
                                from pptx.enum.text import PP_ALIGN  # noqa: F401
                                ph = shape.placeholder_format
                                if ph is not None and ph.idx == 0:
                                    title_text = shape_text
                                    continue  # title included separately
                            except Exception:
                                pass
                        if shape_text:
                            text_parts.append(shape_text)

                    # ---- Images ----
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            img_counter += 1
                            image = shape.image
                            img_bytes = image.blob
                            img_ext = image.ext or "png"

                            img_filename = f"slide_{slide_num}_img_{img_counter}.png"
                            img_path = images_dir / img_filename

                            # Convert to PNG using Pillow for uniformity
                            try:
                                pil_img = Image.open(io.BytesIO(img_bytes))
                                pil_img.save(str(img_path), "PNG")
                            except Exception:
                                # Fall back to raw write
                                img_path.write_bytes(img_bytes)

                            image_records.append({"image_path": str(img_path)})
                        except Exception as exc:
                            logger.warning(
                                "Failed to extract image %d on slide %d of %s: %s",
                                img_counter, slide_num, ppt_path, exc,
                            )

                # ---- Speaker notes ----
                notes_text = ""
                try:
                    notes_slide = slide.notes_slide
                    if notes_slide and notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text.strip()
                except Exception:
                    notes_text = ""

                slides_data.append(
                    {
                        "slide_num": slide_num,
                        "title": title_text,
                        "text": "\n\n".join(text_parts),
                        "images": image_records,
                        "notes": notes_text,
                    }
                )

            logger.info(
                "Extracted %d slides from %s (%d images total)",
                len(slides_data),
                ppt_path,
                sum(len(s["images"]) for s in slides_data),
            )
            return slides_data

        loop = asyncio.get_event_loop()
        try:
            slides = await loop.run_in_executor(None, _extract)
        except PPTIngestionError:
            raise
        except Exception as exc:
            raise PPTIngestionError(
                f"Unexpected error extracting {ppt_path}: {exc}"
            ) from exc

        return slides

    # ------------------------------------------------------------------ #
    # Metadata
    # ------------------------------------------------------------------ #

    def get_metadata(self, ppt_path: Path) -> dict:
        """
        Return presentation-level metadata.

        Returns
        -------
        dict with keys: num_slides, title, author, file_size
        """
        ppt_path = Path(ppt_path)
        if not ppt_path.exists():
            raise PPTIngestionError(f"PPTX file not found: {ppt_path}")

        try:
            prs = Presentation(str(ppt_path))
        except Exception as exc:
            raise PPTIngestionError(
                f"python-pptx could not open {ppt_path}: {exc}"
            ) from exc

        num_slides = len(prs.slides)

        # Core properties are optional
        core = prs.core_properties
        title = getattr(core, "title", "") or ""
        author = getattr(core, "author", "") or ""

        return {
            "num_slides": num_slides,
            "title": title,
            "author": author,
            "file_size": ppt_path.stat().st_size,
        }
